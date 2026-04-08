"""Build Fig.2 platform architecture — v7 简洁干净版."""
from pptx import Presentation
from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image as PILImage
import os, math

MEDIA = os.path.dirname(os.path.abspath(__file__))
SLIDE_W, SLIDE_H = Cm(19.3), Cm(12.0)

COLORS = {
    'purple': (RGBColor(0x5C,0x6B,0xC0), RGBColor(0x79,0x86,0xCB)),
    'green':  (RGBColor(0x2E,0x7D,0x32), RGBColor(0x66,0xBB,0x6A)),
    'orange': (RGBColor(0xE6,0x51,0x00), RGBColor(0xFF,0xA7,0x26)),
    'blue':   (RGBColor(0x15,0x65,0xC0), RGBColor(0x42,0xA5,0xF5)),
}
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_corner_radius(shape, radius=10000):
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is None:
        return
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
    for gd in avLst.findall(qn('a:gd')):
        avLst.remove(gd)
    gd = etree.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {radius}')


def add_img(slide, img_file, area_l, area_t, area_w, area_h):
    """Add image centered, aspect-ratio preserved."""
    img_path = os.path.join(MEDIA, img_file)
    with PILImage.open(img_path) as im:
        iw, ih = im.size
    scale = min(int(area_w) / iw, int(area_h) / ih)
    pw = Emu(int(iw * scale))
    ph = Emu(int(ih * scale))
    pl = area_l + (area_w - pw) // 2
    pt_ = area_t + (area_h - ph) // 2
    slide.shapes.add_picture(img_path, pl, pt_, pw, ph)


def make_card(slide, left, top, width, height, color_key, label, img_file,
              header_ratio=0.25, font_size=9, border_w=Cm(0.08)):
    """
    彩色圆角框 + 内部白色矩形覆盖下部 = 彩色顶 + 白色底.
    框的圆角自然露出彩色边，简洁干净。
    """
    fill_c, border_c = COLORS[color_key]
    bw = int(border_w)

    # 1) 彩色圆角矩形（整个卡片）
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill_c
    bg.line.color.rgb = fill_c
    bg.line.width = Pt(1)
    set_corner_radius(bg, 8000)
    bg.text = ''

    # 2) 白色矩形覆盖下部（留 border_w 给左右底边框）
    header_h = int(height * header_ratio)
    white_top = top + header_h
    white_h = height - header_h - bw
    white = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + bw, white_top,
        width - bw * 2, white_h)
    white.fill.solid()
    white.fill.fore_color.rgb = WHITE
    white.line.fill.background()
    white.text = ''

    # 3) 标签文字在彩色区域
    txBox = slide.shapes.add_textbox(left, top, width, header_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 4) 图片居中在白色区域
    pad = Cm(0.15)
    add_img(slide, img_file,
            left + pad, white_top + pad,
            width - pad * 2, white_h - pad * 2)

    return bg


def add_curved_arrow(slide, x1, y1, x2, y2, color_key, arc_cm=0.8):
    fill_c, _ = COLORS[color_key]
    cp1x = x1 + (x2 - x1) * 0.35
    cp2x = x1 + (x2 - x1) * 0.65
    mid_y = (y1 + y2) / 2
    arc_emu = int(Cm(arc_cm))
    cp1y = mid_y - arc_emu
    cp2y = mid_y - arc_emu

    pts = []
    n = 50
    for i in range(n + 1):
        t = i / n
        bx = (1-t)**3*x1 + 3*(1-t)**2*t*cp1x + 3*(1-t)*t**2*cp2x + t**3*x2
        by = (1-t)**3*y1 + 3*(1-t)**2*t*cp1y + 3*(1-t)*t**2*cp2y + t**3*y2
        pts.append((int(bx), int(by)))

    builder = slide.shapes.build_freeform(Emu(pts[0][0]), Emu(pts[0][1]))
    builder.add_line_segments([(Emu(px), Emu(py)) for px, py in pts[1:]])
    shape = builder.convert_to_shape()
    shape.fill.background()
    shape.line.color.rgb = fill_c
    shape.line.width = Pt(2.2)

    ln = shape._element.find('.//' + qn('a:ln'))
    if ln is not None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')
    return shape


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Robot card ──
    rc_l, rc_t = Cm(0.4), Cm(0.4)
    rc_w, rc_h = Cm(6.8), Cm(11.2)
    make_card(slide, rc_l, rc_t, rc_w, rc_h,
              'purple', 'Autonomous UGV Platform',
              'fig_platform_robot.png',
              header_ratio=0.11, font_size=11)

    # ── Sub cards ──
    sc_w, sc_h = Cm(5.6), Cm(3.3)
    sc_l = SLIDE_W - Cm(0.4) - sc_w
    sc_gap = Cm(0.35)
    sc_tops = [
        Cm(0.4),
        Cm(0.4) + sc_h + sc_gap,
        Cm(0.4) + 2 * (sc_h + sc_gap),
    ]
    sub_data = [
        ('(a) Livox Mid-360 LiDAR',  'fig_platform_lidar.png',   'green'),
        ('(b) Jetson AGX Orin',       'fig_platform_jetson.png',  'orange'),
        ('(c) Ackermann Chassis',     'fig_platform_chassis.png', 'blue'),
    ]
    for i, (label, img, ckey) in enumerate(sub_data):
        make_card(slide, sc_l, sc_tops[i], sc_w, sc_h,
                  ckey, label, img,
                  header_ratio=0.26, font_size=10)

    # ── Curved arrows ──
    arr_x1 = int(rc_l + rc_w)
    arr_x2 = int(sc_l)
    robot_ys = [
        int(rc_t + Cm(2.5)),
        int(rc_t + rc_h / 2),
        int(rc_t + rc_h - Cm(2.5)),
    ]
    sub_ys = [int(sc_tops[i] + sc_h / 2) for i in range(3)]
    arcs = [0.7, 0.1, -0.7]
    ckeys = ['green', 'orange', 'blue']
    for i in range(3):
        add_curved_arrow(slide, arr_x1, robot_ys[i],
                         arr_x2, sub_ys[i], ckeys[i], arc_cm=arcs[i])

    out = os.path.join(MEDIA, 'fig2_platform_v2.pptx')
    prs.save(out)
    print(f'Saved: {out}')


if __name__ == '__main__':
    build()
